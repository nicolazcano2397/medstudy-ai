import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document as DocxDocument

def extract_pdf_text(file_path: str) -> str:
    doc = fitz.open(file_path)
    pages = [page.get_text() for page in doc]
    doc.close()
    return "\n".join(pages).strip()

def extract_ppt_text(file_path: str) -> str:
    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_content = f"\n--- Diapositiva {i} ---\n"
        slide_content += "\n".join(
            shape.text for shape in slide.shapes if hasattr(shape, "text")
        )
        slides_text.append(slide_content)
    return "\n".join(slides_text).strip()

def extract_docx_text(file_path: str) -> str:
    doc = DocxDocument(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs).strip()

def process_uploaded_file(file_path: str, file_type: str) -> str:
    extractors = {
        "pdf": extract_pdf_text,
        "ppt": extract_ppt_text,
        "pptx": extract_ppt_text,
        "doc": extract_docx_text,
        "docx": extract_docx_text,
    }
    if file_type not in extractors:
        raise ValueError(f"Tipo no soportado: {file_type}")
    return extractors[file_type](file_path)
